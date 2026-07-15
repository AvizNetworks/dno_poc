#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

def C(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG   = C("0b0f17"); PANEL=C("141b27"); BD=C("27313f"); FG=C("e9eef5"); MUT=C("93a1b3")
NV   = C("9b57de"); CODE=C("070a10"); CODEFG=C("cdd9e5")   # NV repurposed as Aviz purple accent
LOGO = "/home/ilan/work/dno_poc/hbn_latest/dpf/docs/aviz_logo.png"
CLUSTER=C("a371f7"); NODE=C("3fb950"); POD=C("58a6ff"); CONT=C("f0883e"); IFC=C("39c5cf"); DP=C("e3b341")
SANS="Calibri"; MONO="Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
ML = Inches(0.9)           # left margin
CW = Inches(11.5)          # content width
BLANK = prs.slide_layouts[6]

def new_slide(eyebrow, title, title_color=FG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    # Aviz logo (top-right) + footer
    s.shapes.add_picture(LOGO, SW-Inches(2.05), Inches(0.42), width=Inches(1.55))
    ft = s.shapes.add_textbox(ML, SH-Inches(0.52), CW, Inches(0.35)).text_frame
    ft.word_wrap = False
    fp = ft.paragraphs[0]; fr = fp.add_run(); fr.text = "Aviz Networks   ·   DPF + HBN on BlueField-3"
    fr.font.size = Pt(9); fr.font.color.rgb = MUT; fr.font.name = SANS
    # accent tick
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, Inches(0.52), Inches(0.06), Inches(0.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = NV; bar.line.fill.background()
    # eyebrow
    eb = s.shapes.add_textbox(ML+Inches(0.18), Inches(0.45), CW, Inches(0.4)).text_frame
    eb.word_wrap = True
    p = eb.paragraphs[0]; r = p.add_run(); r.text = eyebrow.upper()
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NV; r.font.name = SANS
    # title
    tt = s.shapes.add_textbox(ML, Inches(0.85), CW, Inches(1.05)).text_frame
    tt.word_wrap = True
    p = tt.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = title_color; r.font.name = SANS
    return s

def textbox(s, left, top, w, h):
    tf = s.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, runs, size=15, color=FG, bold=False, space_after=8, bullet=None, first=False, font=SANS):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after); p.space_before = Pt(0)
    if bullet:
        b = p.add_run(); b.text = "▪  "; b.font.size = Pt(size); b.font.color.rgb = bullet; b.font.bold = True; b.font.name = font
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for tup in runs:
        txt, col, bd = (tup + (color,bold))[:3] if len(tup)<3 else tup
        r = p.add_run(); r.text = txt; r.font.size = Pt(size); r.font.color.rgb = col; r.font.bold = bd; r.font.name = font
    return p

def lead(s, text, top, size=17, color=MUT, w=CW, left=ML):
    tf = textbox(s, left, top, w, Inches(1.0))
    add_para(tf, text, size=size, color=color, first=True)
    return tf

def codebox(s, lines, top, height, width=CW, left=ML, size=12.5):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = CODE
    box.line.color.rgb = NV; box.line.width = Pt(0.75)
    # left accent feel via line only; set radius small
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for ln in lines:
        # ln is list of (text,color) segments or str
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(1); p.line_spacing = 1.15
        segs = ln if isinstance(ln, list) else [(ln, CODEFG)]
        if not segs:
            r = p.add_run(); r.text = " "; r.font.size = Pt(size); r.font.name = MONO
            continue
        for txt, col in segs:
            r = p.add_run(); r.text = txt; r.font.size = Pt(size); r.font.color.rgb = col; r.font.name = MONO
    return box

def note(s, text, top, w=CW, left=ML):
    tf = textbox(s, left, top, w, Inches(0.7))
    add_para(tf, text, size=14, color=MUT, first=True)
    return tf

def card(s, left, top, w, h, title, body, tcolor):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = PANEL; r.line.color.rgb = BD; r.line.width = Pt(1)
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_left=Inches(0.14); tf.margin_right=Inches(0.12); tf.margin_top=Inches(0.1); tf.margin_bottom=Inches(0.08)
    p = tf.paragraphs[0]; run = p.add_run(); run.text = title
    run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = tcolor; run.font.name = SANS
    p.space_after = Pt(4)
    p2 = tf.add_paragraph(); run = p2.add_run(); run.text = body
    run.font.size = Pt(12); run.font.color.rgb = MUT; run.font.name = SANS
    return r

def style_table(tbl, col_widths, header=True):
    # remove banding style by setting fills manually
    for ci,w in enumerate(col_widths):
        tbl.columns[ci].width = w
    for ri, row in enumerate(tbl.rows):
        for ci, cell in enumerate(row.cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = C("1a2230") if (header and ri==0) else PANEL
            cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.06)
            cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = SANS

def set_cell(cell, text, size=12, color=FG, bold=False, mono=False):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = MONO if mono else SANS

# ----------------------------------------------------------------------------
# 1 TITLE
s = prs.slides.add_slide(BLANK)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
s.shapes.add_picture(LOGO, ML, Inches(0.75), width=Inches(3.0))
tt = textbox(s, ML, Inches(2.2), CW, Inches(1.8))
add_para(tt, [("DPF + HBN on ", FG, True), ("BlueField-3", NV, True)], size=48, bold=True, first=True)
lead(s, "Provisioning and operating DPUs the Kubernetes way — flash, join, deploy HBN, and drive FRR from one control point.", Inches(3.7), size=20, color=MUT)
ld = textbox(s, ML, Inches(5.1), CW, Inches(0.6))
add_para(ld, [("Multi-DPU validated (S4 + S2 simultaneously)  ·  DOCA v25.10.1  ·  run it with ", MUT, False), ("dpf/QUICKSTART.md", FG, True)], size=15, first=True)

# 2 WHAT IS DPF
s = new_slide("1 · The idea  / 15", "What is DPF?")
lead(s, [("DPF = DOCA Platform Framework", FG, True), ("  (NVIDIA). Kubernetes-native lifecycle management for BlueField-3 DPUs.", MUT, False)], Inches(1.95), size=18)
tf = textbox(s, ML, Inches(2.9), CW, Inches(3.6))
add_para(tf, [("You ", FG,False),("declare the desired state", FG,True),(" (which OS image, which cluster, which workloads) as Kubernetes objects.", FG,False)], bullet=NV, first=True, size=17, space_after=14)
add_para(tf, [("DPF ", FG,False),("makes it happen", FG,True),(": flashes the BFB → boots the BF3 → joins it to a cluster → deploys services (HBN) onto it.", FG,False)], bullet=NV, size=17, space_after=14)
add_para(tf, [("One control point for many DPUs", FG,True),(" — no SSH to each box. The same kubectl drives 1 or 100.", FG,False)], bullet=NV, size=17)

# 3 K8S
s = new_slide("2 · Foundation  / 15", "Kubernetes in 60 seconds")
layers = [("CLUSTER", CLUSTER, "control plane (API + scheduler + etcd) that runs & tracks everything"),
          ("NODE", NODE, "a machine running kubelet that workloads land on"),
          ("POD", POD, "1+ containers sharing one network identity"),
          ("CONTAINER", CONT, "a process isolated by Linux namespaces + cgroups")]
y = Inches(2.0)
for i,(name,col,desc) in enumerate(layers):
    left = ML + Inches(0.5*i)
    w = Inches(6.0) - Inches(0.5*i)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, w, Inches(1.05))
    r.fill.solid(); r.fill.fore_color.rgb = PANEL; r.line.color.rgb = col; r.line.width = Pt(2)
    tf=r.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.12); tf.margin_top=Inches(0.06); tf.margin_bottom=Inches(0.04)
    p=tf.paragraphs[0]; run=p.add_run(); run.text=name; run.font.size=Pt(13); run.font.bold=True; run.font.color.rgb=col; run.font.name=SANS
    p2=tf.add_paragraph(); run=p2.add_run(); run.text=desc; run.font.size=Pt(11); run.font.color.rgb=MUT; run.font.name=SANS
    y = y + Inches(1.18)
tf = textbox(s, Inches(7.4), Inches(2.0), Inches(5.0), Inches(3.6))
add_para(tf, [("Declarative: ", FG,True),("you describe the end state; controllers reconcile to it.", FG,False)], bullet=NV, first=True, size=15, space_after=12)
add_para(tf, [("kubectl ", FG,True),("is the remote control — talks to the cluster API, not the machines.", FG,False)], bullet=NV, size=15, space_after=12)
add_para(tf, [("DaemonSet", FG,True),(" = “one pod on every node” → deploy once, runs fleet-wide.", FG,False)], bullet=NV, size=15, space_after=12)
add_para(tf, "DPF builds on these primitives — the BF3 becomes a node, HBN becomes a pod.", size=14, color=MUT, space_after=0)

# 4 BUILDING BLOCKS
s = new_slide("3 · NVIDIA's building blocks  / 15", "What NVIDIA built for DPF")
lead(s, "The DPF Operator + sub-controllers (provisioning · dpuservice · kamaji-cm · servicechainset), driven by custom resources you declare:", Inches(1.95), size=16)
cards = [("DPFOperatorConfig", CLUSTER, "Bootstraps the whole system on the DPF VM."),
         ("BFB", DP, "The DPU OS image (DOCA bf-bundle) served to the BF3."),
         ("DPUFlavor", CONT, "How to flash: NIC firmware, hugepages, SubFunctions, config files."),
         ("DPUCluster", CLUSTER, "A per-DPU virtual k8s control plane, run by Kamaji as pods on the DPF VM."),
         ("DPU / DPUNode / DPUDevice", NODE, "The physical BF3 being provisioned and tracked."),
         ("DPUService", POD, "Workloads pushed onto DPUs — CNI, and HBN.")]
cw = Inches(3.7); ch = Inches(1.7); gx = Inches(0.18); gy = Inches(0.22)
x0 = ML; y0 = Inches(2.85)
for i,(t,col,b) in enumerate(cards):
    r = i//3; c = i%3
    card(s, x0 + c*(cw+gx), y0 + r*(ch+gy), cw, ch, t, b, col)

# 5 TOPOLOGY
s = new_slide("4 · Topology  / 15", "Our setup")
topo = [
 [("DPF Operator VM", NV),("  (S5 · 10.4.5.136)            ", CODEFG),("# the control plane", MUT)],
 [("  k3s · DPF Operator v25.10.1 · Kamaji · ArgoCD · bfb-registry", CODEFG)],
 [("        │", CODEFG)],
 [("        │  flash BFB  (rshim via x86 host, or Redfish)", CODEFG)],
 [("        │  Kamaji virtual control plane  ◄── BF3 kubelet joins", CODEFG)],
 [("        ▼", CODEFG)],
 [("BF3 #1 (S4)", NV),("  worker1 · apiserver :6443        ", CODEFG),("# one operator,", MUT)],
 [("BF3 #2 (S2)", NV),("  worker2 · apiserver :6444        ", CODEFG),("# many DPUs", MUT)],
 [("   each BF3 = a Kubernetes WORKER node → runs HBN (FRR) as a pod", CODEFG)],
 [("        │  PCIe", CODEFG)],
 [("        ▼", CODEFG)],
 [("x86 host", NV),("  (10.20.13.226)                       ", CODEFG),("# NOT in k8s", MUT)],
 [("   BF3 rshim (flash) · SR-IOV VFs (vf0..vf7)", CODEFG)],
]
codebox(s, topo, Inches(1.95), Inches(3.7), size=13)
note(s, "Subnet note: in this lab, TCP 10.20.13.x → 10.4.5.x is blocked → a small SSH tunnel bridges the BF3's join. Same-subnet setups skip it.", Inches(5.85))

# 6 WHO DOES WHAT
s = new_slide("5 · Who does what  / 15", "DPF VM ↔ BF3 ↔ host")
rows = [("Box","Role","Runs"),
        ("DPF VM\n10.4.5.136","Control plane / operator","k3s, DPF Operator, Kamaji (DPU control planes), ArgoCD, bfb-registry, kubectl"),
        ("BF3\ns4-dpu","k8s worker node","kubelet + crictl (no kubectl); the doca-hbn pod = FRR/zebra/NVUE; CNI pods"),
        ("x86 host\n10.20.13.226","Not in k8s","BF3 PCIe rshim (OS flash) + SR-IOV VFs vf0..vf7")]
gt = s.shapes.add_table(4, 3, ML, Inches(2.0), CW, Inches(3.0)).table
style_table(gt, [Inches(2.4), Inches(3.0), Inches(6.1)])
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        set_cell(gt.cell(ri,ci), val, size=12, color=(FG if ri==0 else FG), bold=(ri==0 or ci==0))
note(s, "Key point: the x86 host never joins Kubernetes — only the BF3 ARM does.", Inches(5.5))

# 7 BRINGUP
s = new_slide("6 · The bringup  / 15", "How a BF3 comes up")
steps = [("1","Stage BFB","copy the image to the DPF VM"),
         ("2","Run bringup","bringup_dpf.sh --worker worker1 — IPs/creds from config.yaml"),
         ("3","First boot","2 console commands: password + systemctl start dpf-firstboot-kick"),
         ("4","Re-run","DPU goes Ready, HBN DaemonSet deploys"),
         ("5","Host VFs","setup_host_vfs.sh → vf0..vf7")]
sw = Inches(2.18); sh_ = Inches(2.2); gx = Inches(0.15)
x0 = ML; y0 = Inches(2.3)
for i,(n,x,y) in enumerate(steps):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x0+i*(sw+gx), y0, sw, sh_)
    r.fill.solid(); r.fill.fore_color.rgb=PANEL; r.line.color.rgb=BD; r.line.width=Pt(1)
    tf=r.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.12); tf.margin_top=Inches(0.12)
    p=tf.paragraphs[0]; run=p.add_run(); run.text=n; run.font.size=Pt(14); run.font.bold=True; run.font.color.rgb=NV; run.font.name=MONO
    p.space_after=Pt(4)
    p2=tf.add_paragraph(); run=p2.add_run(); run.text=x; run.font.size=Pt(14); run.font.bold=True; run.font.color.rgb=FG; run.font.name=SANS
    p2.space_after=Pt(3)
    p3=tf.add_paragraph(); run=p3.add_run(); run.text=y; run.font.size=Pt(11.5); run.font.color.rgb=MUT; run.font.name=SANS
note(s, "The script is idempotent — run → do the first-boot steps → re-run to finish. Full runbook: dpf/QUICKSTART.md.", Inches(4.9))

# 8 HOST VFs
s = new_slide("7 · Host ports  / 15", "Create the workload ports (host VFs)")
lead(s, "SR-IOV VFs on the x86 host are the ports your VMs / workloads plug into. One command creates and names them.", Inches(1.95), size=17)
codebox(s, [
 [("# on the x86 host, after HBN is up", MUT)],
 [("sudo ./dpf/scripts/setup_host_vfs.sh            ", CODEFG),("# → vf0 .. vf7", MUT)],
 [("sudo ./dpf/scripts/setup_host_vfs.sh --persist  ", CODEFG),("# survive reboot", MUT)],
], Inches(2.9), Inches(1.3), size=13)
tf = textbox(s, ML, Inches(4.5), CW, Inches(2.4))
add_para(tf, [("Auto-detects", FG,True),(" the BlueField-3 PFs by PCI id — no NIC names to look up.", FG,False)], bullet=NV, first=True, size=16, space_after=12)
add_para(tf, [("Names them vf0..vf3 (PF0) and vf4..vf7 (PF1).", FG,False)], bullet=NV, size=16, space_after=12)
add_para(tf, [("Refuses to run with >1 BF3", FG,True),(" on a host (avoids ambiguous naming).", FG,False)], bullet=NV, size=16)

# 9 MENTAL MODEL
s = new_slide("8 · Mental model  / 15", "Think of the BF3 as a router")
diag = [
 [("          ToR / network", CODEFG)],
 [("              │   │", CODEFG)],
 [("            ", CODEFG),("p0  p1", NV),("           ← physical uplinks", MUT)],
 [("      ┌───────┴───┴───────────┐", CODEFG)],
 [("      │   BF3  =  ROUTER (FRR) │", CODEFG)],
 [("      │   uplinks: p0_if p1_if │", CODEFG)],
 [("      │   ports  : pf0vfN_if   │", CODEFG)],
 [("      └───────────┬───────────┘", CODEFG)],
 [("            ", CODEFG),("vf0 … vf7", NV),("          ← workloads attach", MUT)],
 [("        x86 host / VMs & apps", CODEFG)],
]
codebox(s, diag, Inches(2.0), Inches(3.7), width=Inches(6.0), size=12.5)
maprows = [("Host port","Router (FRR)","Role"),
           ("vf0","pf0vf0_if","workload"),("vf1","pf0vf1_if","workload"),
           ("vf2","pf0vf2_if","workload"),("vf3","pf0vf3_if","workload"),
           ("vf4","pf1vf0_if","workload"),("vf5","pf1vf1_if","workload"),
           ("vf6","pf1vf2_if","workload"),("vf7","pf1vf3_if","workload"),
           ("p0 / p1","p0_if / p1_if","uplinks")]
mt = s.shapes.add_table(len(maprows), 3, Inches(8.2), Inches(2.0), Inches(4.3), Inches(4.2)).table
style_table(mt, [Inches(1.4), Inches(1.7), Inches(1.2)])
for ri,row in enumerate(maprows):
    for ci,val in enumerate(row):
        set_cell(mt.cell(ri,ci), val, size=11.5,
                 color=(MUT if ri==0 else (IFC if ci<2 else FG)),
                 bold=(ri==0), mono=(ri>0 and ci<2))

# 10 DEMO UP
s = new_slide("9 · Demo  / 15", "Is everything up?")
codebox(s, [
 [("# on the DPF Operator VM — operator + provisioning state", MUT)],
 [("kubectl get dpfoperatorconfig -n dpf-operator-system        ", CODEFG),("# Ready=True", MUT)],
 [("kubectl get dpu,dpucluster,tenantcontrolplane -n dpf-operator-system", CODEFG)],
 [("./dpf/scripts/fleet_status.sh --frr                          ", CODEFG),("# ALL DPUs: node/HBN/FRR", MUT)],
], Inches(2.0), Inches(1.7), size=13)
codebox(s, [
 [("DPU s4-dpu:        Ready", FG)],
 [("DPUCluster:        True   Ready   kamaji", FG)],
 [("TenantControlPlane:  Ready  10.4.5.136:6443", FG)],
], Inches(4.0), Inches(1.5), size=13)

# 11 DEMO NODE
s = new_slide("10 · Demo  / 15", "The BF3 is a worker node")
codebox(s, [
 [("# point kubectl at the DPU cluster (Kamaji)", MUT)],
 [("D=\"kubectl --kubeconfig /home/dpu-vm/dpu-tc-kubeconfig\"", CODEFG)],
 [("", CODEFG)],
 [("$D get nodes -o wide", CODEFG)],
 [("NAME     STATUS   ROLES    VERSION", FG)],
 [("s4-dpu   Ready    <none>   v1.34.4      ", FG),("# ← the BF3", MUT)],
 [("", CODEFG)],
 [("$D get pods -A | grep -E 'doca-hbn|flannel|coredns'", CODEFG)],
 [("doca-hbn   doca-hbn-xxxxx   1/1   Running   ", FG),("# ← HBN", MUT)],
], Inches(2.0), Inches(3.2), size=12.5)
note(s, "100 DPUs = 100 worker nodes in this view. Add one → its HBN DaemonSet pod schedules automatically.", Inches(5.5))

# 12 DEMO PAYOFF
s = new_slide("11 · Demo — the payoff  / 15", "Reach FRR with kubectl — no SSH")
codebox(s, [
 [("POD=$($D get pod -n doca-hbn -o jsonpath='{.items[0].metadata.name}')", CODEFG)],
 [("", CODEFG)],
 [("# drop into FRR's CLI (zebra/bgpd/staticd answer through vtysh)", MUT)],
 [("$D exec -it -n doca-hbn $POD -- ", CODEFG),("vtysh", NV)],
 [("", CODEFG)],
 [("doca-hbn# show interface brief        ", FG),("# p0_if, pf0vf0_if …", MUT)],
 [("doca-hbn# show ip route               ", FG),("# zebra's RIB", MUT)],
 [("", CODEFG)],
 [("# one-shot, no shell:", MUT)],
 [("$D exec -n doca-hbn $POD -- vtysh -c \"show interface brief\"", CODEFG)],
], Inches(2.0), Inches(3.5), size=12.5)
note(s, "kubectl exec rides the cluster → BF3 kubelet → into the HBN container. Same command for any DPU in the fleet.", Inches(5.8))

# 13 CAVEATS
s = new_slide("12 · Current scope & caveats  / 15", "What works today — and what doesn't yet")
tf = textbox(s, ML, Inches(2.1), CW, Inches(4.5))
add_para(tf, [("Multi-DPU works", FG,True),(" — one operator manages many BF3s (S4 :6443 + S2 :6444 validated simultaneously; unique apiserver_port per worker in config.yaml).", FG,False)], bullet=NV, first=True, size=16, space_after=14)
add_para(tf, [("Flashing is via rshim, not Redfish. ", FG,True),("The OS flash goes through the x86 host's PCIe rshim (--rshim-install), so the x86 host must be reachable with a working rshim. Redfish works only for a genuine BFB version change; S4 was already on the target version, so rshim was required.", FG,False)], bullet=NV, size=16, space_after=14)
add_para(tf, [("Fresh-flash first boot = 2 console commands", FG,True),(" — password + systemctl start dpf-firstboot-kick; later boots hands-off.", FG,False)], bullet=NV, size=16, space_after=14)
add_para(tf, [("Data plane is CPU-routed today", FG,True),(" — fine for config/feature work; eswitch offload arrives with the DPUService migration (roadmap).", FG,False)], bullet=NV, size=16)

# 14 WORKFLOW
s = new_slide("13 · Your workflow  / 15", "Day-to-day for the team")
wf = [("Provision", "bringup_dpf.sh → BF3 flashed, joined, HBN running. Follow QUICKSTART.md."),
      ("Operate FRR", "Your UI / NVUE / REST manages routing on the FRR interfaces — the primary path."),
      ("Debug", "kubectl exec → vtysh, pod logs, status_dpf.sh. See CHEATSHEET.md.")]
cw=Inches(3.7); ch=Inches(1.9); gx=Inches(0.18)
for i,(t,b) in enumerate(wf):
    card(s, ML+i*(cw+gx), Inches(2.1), cw, ch, t, b, NV)
tf = textbox(s, ML, Inches(4.5), CW, Inches(2.0))
add_para(tf, [("Fleet model: ", FG,True),("one DPF VM, many BF3 nodes — provision and inspect them all from one place.", FG,False)], bullet=NV, first=True, size=16, space_after=12)
add_para(tf, [("First bringup of a new server: ", FG,True),("do it with support on hand (BMC behavior varies per box).", FG,False)], bullet=NV, size=16)

# 15 REFERENCES
s = new_slide("14 · References  / 15", "Where everything lives")
refs = [("Doc","Use it for"),
        ("dpf/QUICKSTART.md","Run it — linear bringup steps + first-boot"),
        ("dpf/CHEATSHEET.md","Operate it — pods/nodes/HBN/vtysh"),
        ("dpf/README.md","Understand it — architecture + every known issue")]
rt = s.shapes.add_table(4, 2, ML, Inches(2.1), Inches(9.5), Inches(2.6)).table
style_table(rt, [Inches(3.6), Inches(5.9)])
for ri,row in enumerate(refs):
    for ci,val in enumerate(row):
        set_cell(rt.cell(ri,ci), val, size=13, color=(MUT if ri==0 else (IFC if ci==0 else FG)),
                 bold=(ri==0), mono=(ri>0 and ci==0))
tf = textbox(s, ML, Inches(5.1), CW, Inches(0.7))
add_para(tf, "Questions? → live demo.", size=18, color=NV, bold=True, first=True)

out = "/home/ilan/work/dno_poc/hbn_latest/dpf/docs/dpf-hbn-deck-aviz.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
